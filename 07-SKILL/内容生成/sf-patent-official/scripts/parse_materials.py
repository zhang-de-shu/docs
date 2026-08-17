#!/usr/bin/env python3
"""材料解析脚本 — 阶段 1 材料读取解析

目的：把 Office 材料转成 Markdown，避免只读纯文本漏掉内嵌内容。

输入：.docx / .pptx 文件路径（支持批量）
输出：同名 .md 文件（保留标题层级、表格、列表；图片导出到 <name>_imgs/），返回 md 路径

实现要点：
- docx：优先 mammoth 转 HTML → 转 Markdown（表格保结构）；mammoth 缺失时
  降级为 python-docx 直接抽取（标题样式/列表/表格）
- pptx：python-pptx 逐页抽取标题/正文/表格/备注，标注页码
- 图片：docx 解包 word/media/*；pptx 取图片 shape，均导出到 <name>_imgs/
- 转换失败的单文件不中断，返回 "ERROR: ..." 并在批量报告里标注

依赖：python-pptx、python-docx；mammoth（可选，提升 docx 保真度）
"""

import os
import re
import sys
import zipfile
from html.parser import HTMLParser


class _HtmlToMd(HTMLParser):
    """mammoth 输出的简单 HTML → Markdown。"""

    BLOCK = {"p", "br", "li", "tr", "table"}

    def __init__(self):
        super().__init__()
        self.out = []
        self._cell = None   # 非 None 表示在 td/th 内
        self._row = None
        self._rows = []
        self._ol_stack = []

    def handle_starttag(self, tag, attrs):
        if re.fullmatch(r"h[1-6]", tag):
            self.out.append("\n\n" + "#" * int(tag[1]) + " ")
        elif tag == "p":
            self.out.append("\n\n")
        elif tag == "br":
            self.out.append("\n")
        elif tag == "li":
            if self._ol_stack and self._ol_stack[-1] is not None:
                self._ol_stack[-1] += 1
                self.out.append(f"\n{self._ol_stack[-1]}. ")
            else:
                self.out.append("\n- ")
        elif tag == "ul":
            self._ol_stack.append(None)
        elif tag == "ol":
            self._ol_stack.append(0)
        elif tag in ("td", "th"):
            self._cell = []
        elif tag == "tr":
            self._row = []
        elif tag in ("strong", "b"):
            self._emit("**")
        elif tag in ("em", "i"):
            self._emit("*")
        elif tag == "img":
            src = dict(attrs).get("src", "")
            self._emit("[内嵌图片]" if src.startswith("data:")
                       else f"![img]({os.path.basename(src)})")

    def handle_endtag(self, tag):
        if tag == "ul" or tag == "ol":
            if self._ol_stack:
                self._ol_stack.pop()
        elif tag in ("td", "th"):
            if self._row is not None and self._cell is not None:
                self._row.append("".join(self._cell).strip().replace("|", "\\|"))
            self._cell = None
        elif tag == "tr":
            if self._row is not None:
                self._rows.append(self._row)
            self._row = None
        elif tag == "table":
            self._flush_table()
        elif tag in ("strong", "b"):
            self._emit("**")
        elif tag in ("em", "i"):
            self._emit("*")

    def handle_data(self, data):
        self._emit(data)

    def _emit(self, s: str):
        if self._cell is not None:
            self._cell.append(s)
        else:
            self.out.append(s)

    def _flush_table(self):
        if not self._rows:
            self._rows = []
            return
        n = max(len(r) for r in self._rows)
        lines = []
        for ri, row in enumerate(self._rows):
            cells = [row[i] if i < len(row) else "" for i in range(n)]
            lines.append("| " + " | ".join(cells) + " |")
            if ri == 0:
                lines.append("| " + " | ".join("---" for _ in range(n)) + " |")
        self.out.append("\n\n" + "\n".join(lines) + "\n\n")
        self._rows = []


def _html_to_md(html: str) -> str:
    p = _HtmlToMd()
    p.feed(html)
    p.close()
    return re.sub(r"\n{3,}", "\n\n", "".join(p.out)).strip() + "\n"


def _export_docx_images(docx_path: str, img_dir: str) -> list:
    names = []
    with zipfile.ZipFile(docx_path) as z:
        for n in z.namelist():
            fn = os.path.basename(n)
            if n.startswith("word/media/") and fn:
                os.makedirs(img_dir, exist_ok=True)
                with open(os.path.join(img_dir, fn), "wb") as f:
                    f.write(z.read(n))
                names.append(fn)
    return names


def _docx_via_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    out = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        style = (p.style.name or "").lower()
        if style.startswith("heading"):
            try:
                level = int(style.replace("heading", "").strip())
            except ValueError:
                level = 1
            out.append("#" * min(level, 6) + " " + t)
        elif style.startswith("list"):
            out.append("- " + t)
        else:
            out.append(t)
    for ti, table in enumerate(doc.tables):
        out.append(f"[表格{ti + 1}]")
        for row in table.rows:
            out.append("| " + " | ".join(c.text.replace("\n", " ").strip()
                                         for c in row.cells) + " |")
    return "\n\n".join(out) + "\n"


def _convert_docx(path: str) -> str:
    try:
        import mammoth
        with open(path, "rb") as f:
            return _html_to_md(mammoth.convert_to_html(f).value)
    except ImportError:
        return _docx_via_docx(path)


def _convert_pptx(path: str, img_dir: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## 第{i}页")
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    out.append("| " + " | ".join(c.text.replace("\n", " ").strip()
                                                 for c in row.cells) + " |")
            elif shape.has_text_frame and shape.text_frame.text.strip():
                prefix = "# " if shape == slide.shapes.title else ""
                out.append(prefix + shape.text_frame.text.strip())
            if shape.shape_type == 13:  # 图片
                try:
                    os.makedirs(img_dir, exist_ok=True)
                    fn = f"slide{i}_{shape.shape_id}.{shape.image.ext}"
                    with open(os.path.join(img_dir, fn), "wb") as f:
                        f.write(shape.image.blob)
                except Exception:
                    pass
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            out.append("> 备注：" + slide.notes_slide.notes_text_frame.text.strip())
    return "\n\n".join(out) + "\n"


def convert(path: str) -> str:
    """把 .docx/.pptx 转成同名 .md，返回 md 路径；失败返回 'ERROR: ...'。"""
    try:
        ext = os.path.splitext(path)[1].lower()
        base = os.path.splitext(path)[0]
        img_dir = base + "_imgs"
        if ext == ".docx":
            md = _convert_docx(path)
            _export_docx_images(path, img_dir)
        elif ext == ".pptx":
            md = _convert_pptx(path, img_dir)
        else:
            return f"ERROR: 不支持的类型 {ext}（仅 .docx/.pptx）"
        md_path = base + ".md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(md)
        return md_path
    except Exception as e:  # 单文件失败不中断批量
        return f"ERROR: {path}: {e}"


if __name__ == "__main__":
    for p in sys.argv[1:]:
        print(convert(p))
